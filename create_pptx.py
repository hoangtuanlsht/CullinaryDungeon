
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ─── Color Palette (Dark Fantasy / Dungeon Theme) ───────────────────────────
CLR_BG_DARK    = RGBColor(0x0D, 0x0D, 0x1A)   # Very dark navy
CLR_BG_MID     = RGBColor(0x14, 0x12, 0x2B)   # Deep purple-black
CLR_ACCENT1    = RGBColor(0xE8, 0xA0, 0x1E)   # Golden amber
CLR_ACCENT2    = RGBColor(0xA0, 0x5C, 0xFF)   # Purple glow
CLR_ACCENT3    = RGBColor(0x3D, 0xE8, 0xD0)   # Teal/cyan highlight
CLR_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_TEXT_LIGHT = RGBColor(0xD4, 0xC8, 0xF0)   # Soft lavender-white
CLR_TEXT_GOLD  = RGBColor(0xE8, 0xA0, 0x1E)
CLR_CARD_BG    = RGBColor(0x1E, 0x1A, 0x3A)   # Card background
CLR_DIVIDER    = RGBColor(0xA0, 0x5C, 0xFF)   # Divider lines

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)

# ─── Helper: fill slide background ──────────────────────────────────────────
def fill_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ─── Helper: add rectangle shape ────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.line.width = Pt(line_width) if line_width else 0
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# ─── Helper: add text box ───────────────────────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT,
                font_name="Segoe UI", word_wrap=True, line_spacing=None):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox

# ─── Helper: add multi-paragraph text box ───────────────────────────────────
def add_textbox_multi(slide, lines, left, top, width, height,
                       font_size=18, bold_first=False,
                       color=None, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
    return txBox

# ─── Helper: decorative horizontal line ─────────────────────────────────────
def add_hline(slide, left, top, width, color=CLR_ACCENT1, height_px=3):
    shape = add_rect(slide, left, top, width, Inches(height_px/96),
                     fill_color=color)
    return shape

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TRANG BÌA (Cover)
# ════════════════════════════════════════════════════════════════════════════
def slide_cover(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_DARK)

    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=CLR_ACCENT1)
    # Bottom accent bar
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_color=CLR_ACCENT2)

    # Left glow panel
    add_rect(slide, 0, 0, Inches(0.6), SLIDE_H, fill_color=RGBColor(0x1A, 0x10, 0x35))
    # Right glow panel
    add_rect(slide, SLIDE_W - Inches(0.6), 0, Inches(0.6), SLIDE_H, fill_color=RGBColor(0x1A, 0x10, 0x35))

    # Decorative diamond lines
    add_rect(slide, Inches(0.8), Inches(0.3), Inches(11.73), Inches(0.04), fill_color=CLR_ACCENT2)

    # School label
    add_textbox(slide, "ĐẠI HỌC THỦY LỢI — KHOA CÔNG NGHỆ THÔNG TIN",
                Inches(0.8), Inches(0.45), Inches(11.73), Inches(0.5),
                font_size=11, color=RGBColor(0x90, 0x80, 0xC0), align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Title line 1
    add_textbox(slide, "THIẾT KẾ VÀ PHÁT TRIỂN",
                Inches(0.8), Inches(1.1), Inches(11.73), Inches(1.0),
                font_size=36, bold=True, color=CLR_TEXT_WHITE, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Title line 2 - Game name highlight
    add_textbox(slide, "TRÒ CHƠI 2D KHÁM PHÁ HẦM NGỤC KẾT HỢP NẤU ĂN",
                Inches(0.8), Inches(1.9), Inches(11.73), Inches(0.9),
                font_size=26, bold=True, color=CLR_ACCENT3, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Game title box
    add_rect(slide, Inches(3.5), Inches(2.9), Inches(6.33), Inches(0.9),
             fill_color=RGBColor(0x2A, 0x10, 0x55), line_color=CLR_ACCENT1, line_width=2)
    add_textbox(slide, "✦  CULINARIA ABYSS  ✦",
                Inches(3.5), Inches(2.95), Inches(6.33), Inches(0.9),
                font_size=30, bold=True, color=CLR_ACCENT1, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Divider
    add_rect(slide, Inches(1.5), Inches(4.0), Inches(10.33), Inches(0.03), fill_color=CLR_ACCENT2)

    # Student info
    info_lines = [
        ("👤  Sinh viên:  HOÀNG MINH TUẤN  —  MSSV: 2251061914",  CLR_TEXT_WHITE, 15, True),
        ("🏛️  Lớp: 64CNTT3   |   Ngành: Công nghệ Phần mềm",       CLR_TEXT_LIGHT, 13, False),
        ("👨‍🏫  GVHD: ThS. Trương Xuân Nam",                          CLR_TEXT_LIGHT, 13, False),
        ("📅  Hà Nội, tháng 7 năm 2026",                             RGBColor(0x80, 0x70, 0xB0), 12, False),
    ]
    y = Inches(4.15)
    for text, color, fsize, bold in info_lines:
        add_textbox(slide, text, Inches(0.8), y, Inches(11.73), Inches(0.45),
                    font_size=fsize, bold=bold, color=color, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        y += Inches(0.4)

    # Unity badge
    add_rect(slide, Inches(5.8), Inches(6.5), Inches(1.73), Inches(0.35),
             fill_color=RGBColor(0x22, 0x22, 0x22), line_color=CLR_ACCENT3, line_width=1)
    add_textbox(slide, "🎮 Unity Engine  |  C#",
                Inches(5.8), Inches(6.52), Inches(1.73), Inches(0.35),
                font_size=9, color=CLR_ACCENT3, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MỞ ĐẦU & BỐI CẢNH
# ════════════════════════════════════════════════════════════════════════════
def slide_intro(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_MID)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT2)

    # Slide number badge
    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
             fill_color=CLR_ACCENT2)
    add_textbox(slide, "02", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Header
    add_textbox(slide, "🔍  MỞ ĐẦU & BỐI CẢNH",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT1)

    # Card 1 - Đặt vấn đề
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(5.9), Inches(5.8),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT2, line_width=1)
    add_textbox(slide, "⚔️  ĐẶT VẤN ĐỀ",
                Inches(0.4), Inches(1.15), Inches(5.7), Inches(0.5),
                font_size=15, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(1.65), Inches(5.7), Inches(0.03), fill_color=CLR_ACCENT2)

    problems = [
        "📈 Thị trường game toàn cầu tăng trưởng mạnh",
        "    → 2023: ~$200 tỷ USD doanh thu",
        "",
        "🎮 Dungeon Crawler là thể loại game có",
        "    lịch sử lâu đời nhưng luôn được",
        "    người chơi yêu thích",
        "",
        "🍳 Kết hợp nấu ăn + khám phá hầm ngục",
        "    → Ý tưởng mới mẻ, chưa phổ biến",
        "",
        "🛠️ Thực hành kỹ năng lập trình Unity/C#",
        "    và áp dụng Design Patterns vào",
        "    dự án thực tế",
    ]
    y = Inches(1.8)
    for line in problems:
        if line:
            add_textbox(slide, line, Inches(0.45), y, Inches(5.6), Inches(0.35),
                        font_size=11.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        y += Inches(0.33)

    # Card 2 - Mục tiêu
    add_rect(slide, Inches(6.5), Inches(1.05), Inches(6.5), Inches(5.8),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT3, line_width=1)
    add_textbox(slide, "🎯  MỤC TIÊU ĐỒ ÁN",
                Inches(6.6), Inches(1.15), Inches(6.3), Inches(0.5),
                font_size=15, bold=True, color=CLR_ACCENT3, font_name="Segoe UI")
    add_rect(slide, Inches(6.6), Inches(1.65), Inches(6.3), Inches(0.03), fill_color=CLR_ACCENT3)

    goals = [
        ("✅", "Xây dựng game 2D Action RPG",     "hoàn chỉnh trên Unity Engine"),
        ("✅", "Kết hợp hệ thống Nấu ăn/Chế tạo", "vào vòng lặp Dungeon Crawler"),
        ("✅", "Áp dụng Design Patterns",           "(Singleton, FSM, Observer...)"),
        ("✅", "6 tầng hầm ngục + Sảnh hội",        "với độ khó tăng dần"),
        ("✅", "13+ quái vật, 3 Boss",               "với AI Finite State Machine"),
        ("✅", "Hệ thống Inventory 30 ô",            "& Shop giao dịch vật phẩm"),
    ]
    y = Inches(1.85)
    for icon, line1, line2 in goals:
        add_textbox(slide, f"{icon}  {line1}",
                    Inches(6.6), y, Inches(6.2), Inches(0.3),
                    font_size=12, bold=True, color=CLR_TEXT_WHITE, font_name="Segoe UI")
        add_textbox(slide, f"     {line2}",
                    Inches(6.6), y + Inches(0.28), Inches(6.2), Inches(0.28),
                    font_size=10.5, color=RGBColor(0xA0, 0x90, 0xD0), font_name="Segoe UI")
        y += Inches(0.75)

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CÔNG NGHỆ & CÔNG CỤ
# ════════════════════════════════════════════════════════════════════════════
def slide_tech(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT2)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT1)
    add_textbox(slide, "03", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_BG_DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, "⚙️  CÔNG NGHỆ & CÔNG CỤ SỬ DỤNG",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT2, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT2)

    # Tech cards row 1
    techs_row1 = [
        ("🎮", "Unity Engine",    "2D Game Engine\nGame Loop, Physics\nTileMap, Cinemachine", CLR_ACCENT1),
        ("💻", "C# (C-Sharp)",    "Ngôn ngữ lập trình\nOOP, Events\nCoroutine, LINQ",        CLR_ACCENT2),
        ("🗃️", "ScriptableObject","Lưu trữ dữ liệu\nItemData, Recipes\nNPC Dialogs",         CLR_ACCENT3),
        ("🔥", "Firebase",        "Tài liệu tham khảo\nCloud Firestore\nRealtimeDatabase",    RGBColor(0xFF,0x6B,0x35)),
    ]
    x = Inches(0.3)
    for icon, name, desc, color in techs_row1:
        add_rect(slide, x, Inches(1.1), Inches(3.1), Inches(2.5),
                 fill_color=CLR_CARD_BG, line_color=color, line_width=2)
        add_textbox(slide, icon, x + Inches(0.1), Inches(1.2), Inches(3.0), Inches(0.6),
                    font_size=28, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        add_textbox(slide, name, x + Inches(0.1), Inches(1.75), Inches(3.0), Inches(0.4),
                    font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        add_rect(slide, x + Inches(0.5), Inches(2.18), Inches(2.1), Inches(0.03), fill_color=color)
        for j, dline in enumerate(desc.split("\n")):
            add_textbox(slide, dline, x + Inches(0.1), Inches(2.25) + Inches(j*0.28),
                        Inches(3.0), Inches(0.28),
                        font_size=10, color=CLR_TEXT_LIGHT, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        x += Inches(3.3)

    # Tech cards row 2
    techs_row2 = [
        ("🗺️", "TileMap",          "Thiết kế Level 2D\nGrid-based map\nCollider tích hợp",    CLR_ACCENT1),
        ("🎨", "Aseprite",          "Vẽ Pixel Art\nAnimation Sprite\nExport spritesheet",     CLR_ACCENT2),
        ("📝", "TextMeshPro",       "Hiển thị UI text\nFont tùy chỉnh\nVietnamese support",    CLR_ACCENT3),
        ("🏊", "Object Pooling",    "Tối ưu Performance\nTái sử dụng object\nKhông Destroy liên tục", RGBColor(0xFF,0x6B,0x35)),
    ]
    x = Inches(0.3)
    for icon, name, desc, color in techs_row2:
        add_rect(slide, x, Inches(3.8), Inches(3.1), Inches(2.5),
                 fill_color=CLR_CARD_BG, line_color=color, line_width=2)
        add_textbox(slide, icon, x + Inches(0.1), Inches(3.9), Inches(3.0), Inches(0.6),
                    font_size=28, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        add_textbox(slide, name, x + Inches(0.1), Inches(4.45), Inches(3.0), Inches(0.4),
                    font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        add_rect(slide, x + Inches(0.5), Inches(4.88), Inches(2.1), Inches(0.03), fill_color=color)
        for j, dline in enumerate(desc.split("\n")):
            add_textbox(slide, dline, x + Inches(0.1), Inches(4.95) + Inches(j*0.28),
                        Inches(3.0), Inches(0.28),
                        font_size=10, color=CLR_TEXT_LIGHT, align=PP_ALIGN.CENTER, font_name="Segoe UI")
        x += Inches(3.3)

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GIỚI THIỆU GAME
# ════════════════════════════════════════════════════════════════════════════
def slide_game_intro(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_MID)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT2)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT3)
    add_textbox(slide, "04", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_BG_DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, "🏰  GIỚI THIỆU GAME  —  CULINARIA ABYSS",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT3, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT3)

    # Genre box
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(12.7), Inches(0.55),
             fill_color=RGBColor(0x2A, 0x10, 0x55), line_color=CLR_ACCENT1, line_width=1)
    add_textbox(slide, "🎮 Thể loại: 2D Action RPG  ·  Dungeon Crawler  ·  Survival  ·  Crafting/Cooking",
                Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.45),
                font_size=13, bold=True, color=CLR_ACCENT1, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Left: Story + World
    add_rect(slide, Inches(0.3), Inches(1.75), Inches(5.9), Inches(5.0),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT2, line_width=1)
    add_textbox(slide, "📖  CỐT TRUYỆN & THẾ GIỚI",
                Inches(0.4), Inches(1.85), Inches(5.7), Inches(0.4),
                font_size=13, bold=True, color=CLR_ACCENT2, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(2.27), Inches(5.7), Inches(0.03), fill_color=CLR_ACCENT2)

    story = [
        "Người chơi vào vai một đầu bếp phiêu lưu,",
        "vượt qua 6 tầng hầm ngục nguy hiểm để",
        "tìm kiếm nguyên liệu quý hiếm.",
        "",
        "🗺️  Bản đồ: 1 Sảnh hội (Guild Hub)",
        "         + 6 tầng hầm ngục (dungeon)",
        "",
        "👾  13+ loại quái vật đa dạng",
        "🐉  3 Boss cực mạnh canh giữ tầng cuối",
        "",
        "⚡  Cơ chế Roguelite: Reset khi chết",
        "    → Tạo áp lực & chiến thuật sinh tồn",
        "",
        "💰  Chỉ lưu số tiền Vàng (Coins) giữa",
        "    các lần chơi",
    ]
    y = Inches(2.4)
    for line in story:
        if line:
            add_textbox(slide, line, Inches(0.45), y, Inches(5.65), Inches(0.3),
                        font_size=11, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        y += Inches(0.3)

    # Right: Gameplay Loop
    add_rect(slide, Inches(6.5), Inches(1.75), Inches(6.5), Inches(5.0),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT1, line_width=1)
    add_textbox(slide, "🔄  VÒNG LẶP GAMEPLAY (Game Loop)",
                Inches(6.6), Inches(1.85), Inches(6.3), Inches(0.4),
                font_size=13, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(6.6), Inches(2.27), Inches(6.3), Inches(0.03), fill_color=CLR_ACCENT1)

    loop_steps = [
        ("1️⃣", "KHÁM PHÁ",    "Đi vào hầm ngục, tìm đường",      CLR_ACCENT3),
        ("2️⃣", "CHIẾN ĐẤU",   "Tiêu diệt quái vật & Boss",         CLR_ACCENT1),
        ("3️⃣", "THU THẬP",    "Nhặt nguyên liệu, vật phẩm",        CLR_ACCENT2),
        ("4️⃣", "NẤU ĂN",      "Chế tạo món ăn hồi máu & buff",     RGBColor(0xFF,0x8C,0x00)),
        ("5️⃣", "SINH TỒN",    "Quản lý kho đồ, mua tại Shop",      CLR_ACCENT3),
        ("6️⃣", "TIẾN SÂU HƠN","Chinh phục tầng tiếp theo",          CLR_ACCENT1),
    ]
    y = Inches(2.4)
    for icon, title, desc, color in loop_steps:
        add_rect(slide, Inches(6.6), y, Inches(6.1), Inches(0.65),
                 fill_color=RGBColor(0x22, 0x1A, 0x40), line_color=color, line_width=1)
        add_textbox(slide, f"{icon}  {title}",
                    Inches(6.7), y + Inches(0.05), Inches(2.0), Inches(0.35),
                    font_size=11, bold=True, color=color, font_name="Segoe UI")
        add_textbox(slide, desc,
                    Inches(8.9), y + Inches(0.13), Inches(3.7), Inches(0.35),
                    font_size=10.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        y += Inches(0.73)

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PHÂN TÍCH HỆ THỐNG
# ════════════════════════════════════════════════════════════════════════════
def slide_analysis(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT3)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT1)
    add_textbox(slide, "05", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_BG_DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, "📊  PHÂN TÍCH YÊU CẦU HỆ THỐNG",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT3, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT3)

    # Use Cases list
    add_textbox(slide, "📋  CÁC USE CASE CHÍNH",
                Inches(0.3), Inches(1.05), Inches(6.0), Inches(0.4),
                font_size=14, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")

    use_cases = [
        ("UC-01", "Di chuyển nhân vật",          "Nhảy, đi lại, leo trèo"),
        ("UC-02", "Tấn công & Phòng thủ",         "Đòn thường, chặn đỡ"),
        ("UC-03", "Quản lý Kho đồ (Inventory)",   "30 ô, kéo thả vật phẩm"),
        ("UC-04", "Nấu ăn & Chế tạo (Crafting)",  "Kiểm tra công thức độc lập"),
        ("UC-05", "Mua bán tại Shop",              "Giao dịch với NPC"),
        ("UC-06", "Tương tác NPC & Hội thoại",     "Dialog system"),
        ("UC-07", "AI Quái vật (FSM)",             "Idle, Patrol, Attack states"),
        ("UC-08", "Dịch chuyển giữa các Scene",    "SceneManager loading"),
        ("UC-09", "Sương mù tầm nhìn (FogOfWar)",  "Overlay shader/mask"),
        ("UC-10", "Tương tác Vật phẩm thế giới",   "Pickup, drop items"),
    ]
    y = Inches(1.5)
    for code, name, desc in use_cases:
        add_rect(slide, Inches(0.3), y, Inches(6.2), Inches(0.42),
                 fill_color=CLR_CARD_BG, line_color=CLR_ACCENT2, line_width=1)
        add_textbox(slide, code,
                    Inches(0.35), y + Inches(0.06), Inches(0.9), Inches(0.3),
                    font_size=9, bold=True, color=CLR_ACCENT2, font_name="Segoe UI")
        add_textbox(slide, name,
                    Inches(1.3), y + Inches(0.06), Inches(2.5), Inches(0.3),
                    font_size=10, bold=True, color=CLR_TEXT_WHITE, font_name="Segoe UI")
        add_textbox(slide, desc,
                    Inches(3.85), y + Inches(0.06), Inches(2.65), Inches(0.3),
                    font_size=9.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        y += Inches(0.47)

    # Actors box
    add_rect(slide, Inches(6.7), Inches(1.05), Inches(6.3), Inches(2.0),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT1, line_width=1)
    add_textbox(slide, "👥  TÁC NHÂN HỆ THỐNG",
                Inches(6.8), Inches(1.12), Inches(6.1), Inches(0.4),
                font_size=13, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(6.8), Inches(1.55), Inches(6.1), Inches(0.03), fill_color=CLR_ACCENT1)
    actors = [
        ("🧑‍🍳", "Người chơi (Player)", "Nhân vật chính, tương tác toàn bộ game"),
        ("👾", "Quái vật / Boss",      "AI tự động theo FSM"),
        ("🤖", "NPC / Shop",           "Bán hàng, thoại, cung cấp nhiệm vụ"),
        ("⚙️", "Hệ thống Game",        "Scene, Spawn, Sound, UI Manager"),
    ]
    ya = Inches(1.65)
    for icon, name, desc in actors:
        add_textbox(slide, f"{icon}  {name}  —  {desc}",
                    Inches(6.8), ya, Inches(6.1), Inches(0.32),
                    font_size=10.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        ya += Inches(0.33)

    # DB / Save
    add_rect(slide, Inches(6.7), Inches(3.2), Inches(6.3), Inches(1.3),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT3, line_width=1)
    add_textbox(slide, "💾  CƠ SỞ DỮ LIỆU LƯU TRỮ",
                Inches(6.8), Inches(3.28), Inches(6.1), Inches(0.35),
                font_size=13, bold=True, color=CLR_ACCENT3, font_name="Segoe UI")
    add_rect(slide, Inches(6.8), Inches(3.65), Inches(6.1), Inches(0.03), fill_color=CLR_ACCENT3)
    db_info = [
        "• Sử dụng PlayerPrefs (Local Storage) của Unity",
        "• Lưu trữ duy nhất: số Tiền Vàng (Coins) qua key 'coin'",
        "• Cơ chế Roguelite: Máu & Kho đồ reset khi chết",
    ]
    yd = Inches(3.72)
    for line in db_info:
        add_textbox(slide, line, Inches(6.8), yd, Inches(6.1), Inches(0.28),
                    font_size=10.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        yd += Inches(0.3)

    # Design Patterns
    add_rect(slide, Inches(6.7), Inches(4.65), Inches(6.3), Inches(2.5),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT2, line_width=1)
    add_textbox(slide, "🏗️  KIẾN TRÚC & DESIGN PATTERNS",
                Inches(6.8), Inches(4.73), Inches(6.1), Inches(0.35),
                font_size=13, bold=True, color=CLR_ACCENT2, font_name="Segoe UI")
    add_rect(slide, Inches(6.8), Inches(5.1), Inches(6.1), Inches(0.03), fill_color=CLR_ACCENT2)
    patterns = [
        ("🔵", "Singleton",         "InventoryManager, SoundManager, UIManager"),
        ("🟣", "FSM",               "AI quái vật: Idle → Patrol → Attack"),
        ("🟡", "Observer/Event",     "SpawnManager lắng nghe SceneLoaded"),
        ("🟢", "ScriptableObject",   "ItemData, CraftingRecipe, NPCDialogs"),
    ]
    yp = Inches(5.17)
    for icon, name, desc in patterns:
        add_textbox(slide, f"{icon}  {name}:", Inches(6.8), yp, Inches(2.3), Inches(0.35),
                    font_size=11, bold=True, color=CLR_TEXT_WHITE, font_name="Segoe UI")
        add_textbox(slide, desc, Inches(9.2), yp, Inches(3.7), Inches(0.35),
                    font_size=10.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        yp += Inches(0.43)

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CÁC HỆ THỐNG GAME CHÍNH
# ════════════════════════════════════════════════════════════════════════════
def slide_systems(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_MID)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT3)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT3)
    add_textbox(slide, "06", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_BG_DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, "⚙️  CÁC HỆ THỐNG GAME CHÍNH",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT1)

    systems = [
        {
            "icon": "⚔️", "title": "HỆ THỐNG CHIẾN ĐẤU",
            "color": CLR_ACCENT1,
            "items": [
                "• Di chuyển: Chạy, Nhảy, Crouch",
                "• Tấn công cận chiến cơ bản",
                "• Phòng thủ (Shield): giảm sát thương",
                "• Hệ thống máu & sát thương",
            ]
        },
        {
            "icon": "🍳", "title": "NẤU ĂN & CHẾ TẠO",
            "color": RGBColor(0xFF,0x8C,0x00),
            "items": [
                "• Kiểm tra công thức độc lập thứ tự",
                "• 30+ vật phẩm nguyên liệu",
                "• Món ăn hồi máu & buff stats",
                "• Tích hợp trực tiếp Gameplay Loop",
            ]
        },
        {
            "icon": "🎒", "title": "KHO ĐỒ (INVENTORY)",
            "color": CLR_ACCENT3,
            "items": [
                "• 30 ô lưu trữ vật phẩm",
                "• Recyclable Scroll Rect (tối ưu render)",
                "• Dùng được, xem thông tin vật phẩm",
                "• Phân loại: Nguyên liệu, Thức ăn, Trang bị",
            ]
        },
        {
            "icon": "👾", "title": "AI QUÁI VẬT (FSM)",
            "color": CLR_ACCENT2,
            "items": [
                "• 3 trạng thái: Idle → Patrol → Attack",
                "• Phát hiện qua EnemySight (trigger)",
                "• Mất dấu → về Idle sau thời gian",
                "• 13+ loại quái  &  3 Boss đặc biệt",
            ]
        },
        {
            "icon": "🌫️", "title": "SƯƠNG MÙ TẦM NHÌN",
            "color": CLR_ACCENT3,
            "items": [
                "• Fog of War overlay shader/mask",
                "• Vùng sáng cập nhật theo Player",
                "• Tạo yếu tố khám phá hấp dẫn",
                "• Tăng tính căng thẳng gameplay",
            ]
        },
        {
            "icon": "🏪", "title": "SHOP & NPC",
            "color": CLR_ACCENT1,
            "items": [
                "• Mua bán vật phẩm bằng Coin",
                "• Dialog system với NPC",
                "• Singleton ShopController",
                "• Spawn Manager qua SceneLoaded",
            ]
        },
    ]

    cols = 3
    card_w = Inches(4.1)
    card_h = Inches(2.55)
    x_start = Inches(0.3)
    y_start = Inches(1.05)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)

    for idx, sys in enumerate(systems):
        row = idx // cols
        col = idx % cols
        x = x_start + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)
        color = sys["color"]

        add_rect(slide, x, y, card_w, card_h,
                 fill_color=CLR_CARD_BG, line_color=color, line_width=2)
        # Top bar
        add_rect(slide, x, y, card_w, Inches(0.07), fill_color=color)
        add_textbox(slide, f"{sys['icon']}  {sys['title']}",
                    x + Inches(0.1), y + Inches(0.1), card_w - Inches(0.2), Inches(0.4),
                    font_size=12, bold=True, color=color, font_name="Segoe UI")
        add_rect(slide, x + Inches(0.1), y + Inches(0.52), card_w - Inches(0.2), Inches(0.03), fill_color=color)

        yi = y + Inches(0.62)
        for item in sys["items"]:
            add_textbox(slide, item, x + Inches(0.12), yi, card_w - Inches(0.22), Inches(0.35),
                        font_size=10, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
            yi += Inches(0.42)

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DEMO GIAO DIỆN
# ════════════════════════════════════════════════════════════════════════════
def slide_demo(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT2)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT2)
    add_textbox(slide, "07", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_TEXT_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "🖥️  DEMO GIAO DIỆN & GAMEPLAY",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT2, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT2)

    screens = [
        ("🏠", "Main Menu",          "Màn hình chính\n(Start / Exit)"),
        ("🏛️", "Guild Hub",          "Sảnh hội\nMua sắm & chuẩn bị"),
        ("⚔️", "Gameplay HUD",       "Thanh máu, Coin\nFog of War"),
        ("🎒", "Inventory & Cooking","Kho đồ 30 ô\n& Bếp nấu ăn"),
        ("💬", "NPC Dialog",         "Hội thoại NPC\n& Shop giao dịch"),
        ("🐉", "Boss Battle",        "Boss Demon Slime\nMinotaurus, Gruz Mother"),
    ]

    x_start = Inches(0.3)
    y_start = Inches(1.05)
    card_w = Inches(4.1)
    card_h = Inches(2.8)
    gap_x = Inches(0.2)

    colors = [CLR_ACCENT1, CLR_ACCENT3, CLR_ACCENT2, RGBColor(0xFF,0x8C,0x00), CLR_ACCENT3, CLR_ACCENT1]

    for row in range(2):
        for col in range(3):
            idx = row * 3 + col
            if idx >= len(screens):
                break
            icon, title, desc = screens[idx]
            color = colors[idx]
            x = x_start + col * (card_w + gap_x)
            y = y_start + row * (card_h + Inches(0.15))

            add_rect(slide, x, y, card_w, card_h,
                     fill_color=CLR_CARD_BG, line_color=color, line_width=2)
            add_rect(slide, x, y, card_w, Inches(0.07), fill_color=color)

            # Icon area (placeholder for screenshot)
            add_rect(slide, x + Inches(0.15), y + Inches(0.15),
                     card_w - Inches(0.3), Inches(1.6),
                     fill_color=RGBColor(0x15, 0x12, 0x25), line_color=color, line_width=1)
            add_textbox(slide, icon,
                        x + Inches(0.15), y + Inches(0.35),
                        card_w - Inches(0.3), Inches(1.0),
                        font_size=44, align=PP_ALIGN.CENTER, font_name="Segoe UI")
            add_textbox(slide, "[Screenshot]",
                        x + Inches(0.15), y + Inches(1.2),
                        card_w - Inches(0.3), Inches(0.4),
                        font_size=9, color=RGBColor(0x50,0x40,0x80),
                        align=PP_ALIGN.CENTER, italic=True, font_name="Segoe UI")

            add_textbox(slide, title,
                        x + Inches(0.1), y + Inches(1.9),
                        card_w - Inches(0.2), Inches(0.4),
                        font_size=13, bold=True, color=color, font_name="Segoe UI")
            for j, dline in enumerate(desc.split("\n")):
                add_textbox(slide, dline,
                            x + Inches(0.1), y + Inches(2.32) + Inches(j*0.23),
                            card_w - Inches(0.2), Inches(0.25),
                            font_size=9.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")

    add_textbox(slide, "💡 Ghi chú: Có thể chèn ảnh chụp màn hình thực tế vào các ô trên",
                Inches(0.4), Inches(7.1), Inches(12.0), Inches(0.3),
                font_size=9, italic=True, color=RGBColor(0x60,0x50,0x90), font_name="Segoe UI")
    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KẾT QUẢ & ĐÁNH GIÁ
# ════════════════════════════════════════════════════════════════════════════
def slide_results(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_MID)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT3)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT1)
    add_textbox(slide, "08", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_BG_DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, "📊  KẾT QUẢ ĐẠT ĐƯỢC & ĐÁNH GIÁ",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT1)

    # Progress bar header
    add_rect(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(0.65),
             fill_color=RGBColor(0x2A, 0x10, 0x55), line_color=CLR_ACCENT1, line_width=1)
    add_textbox(slide, "🎯  Hoàn thành ước tính:  ~85–90%  của các mục tiêu đề ra",
                Inches(0.4), Inches(1.05), Inches(10.0), Inches(0.5),
                font_size=15, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    # Progress bar
    add_rect(slide, Inches(10.8), Inches(1.12), Inches(2.0), Inches(0.35),
             fill_color=RGBColor(0x20,0x20,0x20), line_color=CLR_ACCENT1, line_width=1)
    add_rect(slide, Inches(10.8), Inches(1.12), Inches(1.75), Inches(0.35),
             fill_color=CLR_ACCENT3)
    add_textbox(slide, "87%", Inches(10.8), Inches(1.12), Inches(2.0), Inches(0.35),
                font_size=12, bold=True, color=CLR_BG_DARK, align=PP_ALIGN.CENTER)

    # Left: Achievements
    add_rect(slide, Inches(0.3), Inches(1.8), Inches(6.0), Inches(4.8),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT3, line_width=1)
    add_textbox(slide, "✅  KẾT QUẢ ĐẠT ĐƯỢC",
                Inches(0.4), Inches(1.9), Inches(5.8), Inches(0.4),
                font_size=14, bold=True, color=CLR_ACCENT3, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(2.33), Inches(5.8), Inches(0.03), fill_color=CLR_ACCENT3)

    achievements = [
        "✔ Game 2D Action Platformer mượt mà, build ổn định trên PC",
        "✔ AI quái vật hoạt động đầy đủ: Idle, Patrol, Attack",
        "✔ Hệ thống Crafting/Cooking hoàn chỉnh",
        "✔ Kho đồ Inventory 30 ô + Shop giao dịch",
        "✔ Fog of War (Sương mù tầm nhìn) hoạt động",
        "✔ 6 tầng hầm ngục + 1 Sảnh hội",
        "✔ 13+ quái vật + 3 Boss đa dạng",
        "✔ Object Pooling: tối ưu hiệu năng",
        "✔ NPC Dialog System",
        "✔ Áp dụng đúng các Design Patterns",
    ]
    ya = Inches(2.45)
    for line in achievements:
        add_textbox(slide, line, Inches(0.4), ya, Inches(5.8), Inches(0.38),
                    font_size=11, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        ya += Inches(0.42)

    # Right: Limitations + Future
    add_rect(slide, Inches(6.5), Inches(1.8), Inches(6.5), Inches(2.1),
             fill_color=CLR_CARD_BG, line_color=RGBColor(0xFF,0x60,0x60), line_width=1)
    add_textbox(slide, "⚠️  HẠN CHẾ CÒN TỒN TẠI",
                Inches(6.6), Inches(1.9), Inches(6.3), Inches(0.4),
                font_size=13, bold=True, color=RGBColor(0xFF,0x80,0x80), font_name="Segoe UI")
    add_rect(slide, Inches(6.6), Inches(2.33), Inches(6.3), Inches(0.03), fill_color=RGBColor(0xFF,0x60,0x60))
    limits = [
        "⚠ Save/Load đơn giản — chỉ lưu số Coins",
        "⚠ Chưa có hệ thống kỹ năng (Skill Tree)",
        "⚠ Nhạc nền chưa đầy đủ cho từng tầng",
        "⚠ Chỉ cận chiến, chưa có vũ khí tầm xa",
    ]
    yl = Inches(2.45)
    for line in limits:
        add_textbox(slide, line, Inches(6.6), yl, Inches(6.3), Inches(0.34),
                    font_size=10.5, color=RGBColor(0xFF,0xC0,0xC0), font_name="Segoe UI")
        yl += Inches(0.38)

    add_rect(slide, Inches(6.5), Inches(4.05), Inches(6.5), Inches(2.55),
             fill_color=CLR_CARD_BG, line_color=CLR_ACCENT2, line_width=1)
    add_textbox(slide, "🚀  HƯỚNG PHÁT TRIỂN TƯƠNG LAI",
                Inches(6.6), Inches(4.15), Inches(6.3), Inches(0.4),
                font_size=13, bold=True, color=CLR_ACCENT2, font_name="Segoe UI")
    add_rect(slide, Inches(6.6), Inches(4.58), Inches(6.3), Inches(0.03), fill_color=CLR_ACCENT2)
    futures = [
        "🔮 Tích hợp Firebase / JSON Serialization để lưu trữ đầy đủ",
        "🌳 Phát triển cây kỹ năng (Skill Tree) qua EXP & thức ăn",
        "📋 Hệ thống Nhiệm vụ (Quest System) với NPC",
        "📱 Tối ưu Mobile: Joystick ảo → build Android/iOS",
        "🗺️ Mở rộng vũ khí, áo giáp, tầng hầm sâu hơn",
    ]
    yf = Inches(4.68)
    for line in futures:
        add_textbox(slide, line, Inches(6.6), yf, Inches(6.3), Inches(0.34),
                    font_size=10.5, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
        yf += Inches(0.38)

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KẾT LUẬN
# ════════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT1)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=CLR_ACCENT2)

    add_rect(slide, Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4), fill_color=CLR_ACCENT2)
    add_textbox(slide, "09", Inches(12.5), Inches(0.1), Inches(0.73), Inches(0.4),
                font_size=14, bold=True, color=CLR_TEXT_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "📝  KẾT LUẬN",
                Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.6),
                font_size=24, bold=True, color=CLR_ACCENT1, font_name="Segoe UI")
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.04), fill_color=CLR_ACCENT1)

    # Summary quote box
    add_rect(slide, Inches(0.5), Inches(1.0), Inches(12.33), Inches(1.2),
             fill_color=RGBColor(0x1A, 0x10, 0x35), line_color=CLR_ACCENT1, line_width=2)
    add_rect(slide, Inches(0.5), Inches(1.0), Inches(0.1), Inches(1.2), fill_color=CLR_ACCENT1)
    add_textbox(slide,
                "\"Culinaria Abyss\" là sản phẩm game 2D hoàn chỉnh, thể hiện khả năng\n"
                "vận dụng kiến thức công nghệ phần mềm vào phát triển game thực tế.",
                Inches(0.75), Inches(1.08), Inches(12.0), Inches(0.9),
                font_size=14, italic=True, color=CLR_TEXT_WHITE, font_name="Segoe UI")

    # Key takeaways: 2 columns
    takeaways_left = [
        ("🎓", "Kiến thức áp dụng",
         "Singleton, FSM, Observer, ScriptableObject\nObject Pooling, TileMap, Cinemachine"),
        ("🎮", "Kỹ năng phát triển",
         "Unity C# — Từ khái niệm → Sản phẩm hoàn chỉnh\nPhân tích, thiết kế, kiểm thử game"),
        ("📐", "Phân tích & Thiết kế",
         "Use Case 10 chức năng, Activity & Sequence Diagrams\nKiến trúc hệ thống rõ ràng"),
    ]
    takeaways_right = [
        ("🌟", "Điểm nổi bật",
         "Hệ thống Crafting độc lập thứ tự\nFog of War & AI FSM hoàn chỉnh"),
        ("📦", "Quy mô dự án",
         "6 tầng + 1 Sảnh | 13+ quái | 3 Boss\n30 ô Inventory | 30+ vật phẩm"),
        ("🚀", "Định hướng tương lai",
         "Firebase, Quest System, Mobile Build\nSkill Tree, Multiplayer"),
    ]

    def render_col(items, x_base, y_base):
        y = y_base
        for icon, title, desc in items:
            add_rect(slide, x_base, y, Inches(6.0), Inches(1.45),
                     fill_color=CLR_CARD_BG, line_color=CLR_ACCENT2, line_width=1)
            add_textbox(slide, f"{icon}  {title}",
                        x_base + Inches(0.1), y + Inches(0.1), Inches(5.8), Inches(0.4),
                        font_size=13, bold=True, color=CLR_ACCENT3, font_name="Segoe UI")
            add_rect(slide, x_base + Inches(0.1), y + Inches(0.52), Inches(5.8), Inches(0.03), fill_color=CLR_ACCENT2)
            for j, dline in enumerate(desc.split("\n")):
                add_textbox(slide, dline,
                            x_base + Inches(0.1), y + Inches(0.6) + Inches(j*0.32),
                            Inches(5.8), Inches(0.32),
                            font_size=11, color=CLR_TEXT_LIGHT, font_name="Segoe UI")
            y += Inches(1.6)

    render_col(takeaways_left,  Inches(0.3),  Inches(2.35))
    render_col(takeaways_right, Inches(6.7), Inches(2.35))

    # Bottom bar
    add_rect(slide, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.25),
             fill_color=RGBColor(0x1A, 0x10, 0x35))
    add_textbox(slide, "🏛️ Đại học Thủy Lợi  —  Khoa CNTT  —  Đồ án tốt nghiệp 2026",
                Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.22),
                font_size=9, color=RGBColor(0x70,0x60,0xA0), align=PP_ALIGN.CENTER, font_name="Segoe UI")

    return slide

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TRANG CẢM ƠN (Thank You)
# ════════════════════════════════════════════════════════════════════════════
def slide_thankyou(prs):
    slide = blank_slide(prs)
    fill_bg(slide, CLR_BG_DARK)

    # Decorative top/bottom bars
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.15), fill_color=CLR_ACCENT1)
    add_rect(slide, 0, SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), fill_color=CLR_ACCENT2)
    add_rect(slide, 0, Inches(0.15), SLIDE_W, Inches(0.05), fill_color=CLR_ACCENT2)
    add_rect(slide, 0, SLIDE_H - Inches(0.2), SLIDE_W, Inches(0.05), fill_color=CLR_ACCENT1)

    # Side bars
    add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, fill_color=CLR_ACCENT2)
    add_rect(slide, SLIDE_W - Inches(0.15), 0, Inches(0.15), SLIDE_H, fill_color=CLR_ACCENT1)

    # Main Thank You text
    add_textbox(slide, "✦  XIN TRÂN TRỌNG CẢM ƠN!  ✦",
                Inches(0.3), Inches(1.3), Inches(12.73), Inches(1.2),
                font_size=44, bold=True, color=CLR_ACCENT1,
                align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Game name
    add_rect(slide, Inches(3.0), Inches(2.6), Inches(7.33), Inches(0.75),
             fill_color=RGBColor(0x2A, 0x10, 0x55), line_color=CLR_ACCENT1, line_width=2)
    add_textbox(slide, "🎮  CULINARIA ABYSS  🎮",
                Inches(3.0), Inches(2.68), Inches(7.33), Inches(0.6),
                font_size=26, bold=True, color=CLR_ACCENT3,
                align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Divider
    add_rect(slide, Inches(2.0), Inches(3.5), Inches(9.33), Inches(0.04), fill_color=CLR_ACCENT2)

    # Student info
    info = [
        ("👤  Sinh viên:", "HOÀNG MINH TUẤN  —  MSSV: 2251061914", CLR_ACCENT1, CLR_TEXT_WHITE, 13, 16),
        ("🏛️  Trường:",    "Đại học Thủy Lợi — Khoa Công nghệ Thông tin", CLR_ACCENT3, CLR_TEXT_LIGHT, 12, 13),
        ("📚  Lớp:",       "64CNTT3  |  Ngành: Công nghệ Phần mềm",         CLR_ACCENT3, CLR_TEXT_LIGHT, 12, 13),
        ("👨‍🏫  GVHD:",     "ThS. Trương Xuân Nam",                          CLR_ACCENT2, CLR_TEXT_LIGHT, 12, 13),
    ]
    y = Inches(3.65)
    for label, value, lcolor, vcolor, lfsize, vfsize in info:
        add_textbox(slide, label,
                    Inches(3.2), y, Inches(2.0), Inches(0.45),
                    font_size=lfsize, bold=True, color=lcolor, font_name="Segoe UI")
        add_textbox(slide, value,
                    Inches(5.3), y, Inches(6.5), Inches(0.45),
                    font_size=vfsize, color=vcolor, font_name="Segoe UI")
        y += Inches(0.48)

    # Divider 2
    add_rect(slide, Inches(2.0), Inches(5.65), Inches(9.33), Inches(0.04), fill_color=CLR_ACCENT1)

    # Q&A invite
    add_textbox(slide, "🙋  Kính mời Hội đồng đặt câu hỏi!",
                Inches(0.3), Inches(5.8), Inches(12.73), Inches(0.65),
                font_size=22, bold=True, color=CLR_ACCENT2,
                align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Tech stack footer
    add_rect(slide, Inches(3.5), Inches(6.55), Inches(6.33), Inches(0.4),
             fill_color=RGBColor(0x18,0x15,0x30), line_color=CLR_ACCENT2, line_width=1)
    add_textbox(slide, "🛠️  Unity  ·  C#  ·  TileMap  ·  Cinemachine  ·  Firebase  ·  Aseprite",
                Inches(3.5), Inches(6.58), Inches(6.33), Inches(0.35),
                font_size=9, color=CLR_ACCENT3, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    return slide

# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    print("Creating slides...")
    slide_cover(prs)       # 1
    slide_intro(prs)       # 2
    slide_tech(prs)        # 3
    slide_game_intro(prs)  # 4
    slide_analysis(prs)    # 5
    slide_systems(prs)     # 6
    slide_demo(prs)        # 7
    slide_results(prs)     # 8
    slide_conclusion(prs)  # 9
    slide_thankyou(prs)    # 10

    output_path = "d:/Project/CulinariaDungeon/CulinariaDungeon_BaoVe_DATN.pptx"
    prs.save(output_path)
    print(f"✅ Done! Saved to: {output_path}")

if __name__ == "__main__":
    main()
